"""Build the solution presentation deck.

Generated rather than hand-drawn so the numbers on the evaluation slide come from
evaluation/results/evaluation_report.json instead of being retyped and going stale.

    python scripts/generate_deck.py [-o docs/Invoice-to-Payment-Agent.pptx]
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
REPORT = ROOT / "evaluation" / "results" / "evaluation_report.json"

NAVY = RGBColor(0x0B, 0x25, 0x45)
TEAL = RGBColor(0x00, 0xA1, 0x9A)
AMBER = RGBColor(0xE8, 0x8B, 0x2F)
SLATE = RGBColor(0x33, 0x44, 0x55)
MUTED = RGBColor(0x6B, 0x7A, 0x8B)
PAPER = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD4, 0xDB, 0xE3)

FONT = "Calibri"
W = Inches(13.333)
H = Inches(7.5)


def _text(frame: Any, size: int, color: RGBColor, bold: bool = False) -> None:
    for paragraph in frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = FONT


def box(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str = "",
    fill: RGBColor = WHITE,
    line: RGBColor = BORDER,
    font_size: int = 11,
    font_color: RGBColor = SLATE,
    bold: bool = False,
    shape: Any = MSO_SHAPE.ROUNDED_RECTANGLE,
) -> Any:
    node = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    node.fill.solid()
    node.fill.fore_color.rgb = fill
    node.line.color.rgb = line
    node.line.width = Pt(1)
    node.shadow.inherit = False
    frame = node.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = frame.margin_right = Inches(0.06)
    frame.margin_top = frame.margin_bottom = Inches(0.03)
    frame.text = text
    for paragraph in frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
    _text(frame, font_size, font_color, bold)
    return node


def label(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    size: int = 12,
    color: RGBColor = SLATE,
    bold: bool = False,
    align: Any = PP_ALIGN.LEFT,
) -> Any:
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = tb.text_frame
    frame.word_wrap = True
    frame.text = text
    for paragraph in frame.paragraphs:
        paragraph.alignment = align
    _text(frame, size, color, bold)
    return tb


def arrow(slide: Any, x1: float, y1: float, x2: float, y2: float, color: RGBColor = MUTED) -> None:
    conn = slide.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.75)
    conn.line.end_arrowhead = True  # type: ignore[attr-defined]


def blank(prs: Presentation) -> Any:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def content(prs: Presentation, title: str, kicker: str = "") -> Any:
    slide = blank(prs)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.09))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    bar.shadow.inherit = False

    label(slide, 0.6, 0.42, 12.2, 0.6, title, size=30, color=NAVY, bold=True)
    if kicker:
        label(slide, 0.62, 1.06, 12.2, 0.4, kicker, size=14, color=MUTED)
    return slide


def bullets(
    slide: Any,
    left: float,
    top: float,
    width: float,
    items: list[str],
    size: int = 15,
    gap: float = 0.46,
) -> None:
    for index, item in enumerate(items):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left),
            Inches(top + index * gap + 0.09),
            Inches(0.11),
            Inches(0.11),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = TEAL
        dot.line.fill.background()
        dot.shadow.inherit = False
        label(slide, left + 0.26, top + index * gap - 0.03, width - 0.26, gap, item, size=size)


def metric_tile(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    value: str,
    caption: str,
    accent: RGBColor = TEAL,
) -> None:
    tile = box(slide, left, top, width, height, "", fill=PAPER, line=BORDER)
    tile.text_frame.text = ""
    label(
        slide,
        left,
        top + 0.16,
        width,
        0.6,
        value,
        size=30,
        color=accent,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    label(slide, left, top + 0.78, width, 0.6, caption, size=11, color=MUTED, align=PP_ALIGN.CENTER)


def load_metrics() -> dict[str, Any]:
    if not REPORT.exists():
        return {}
    return json.loads(REPORT.read_text())


def pct(value: Any) -> str:
    try:
        return f"{float(value) * 100:.0f}%"
    except (TypeError, ValueError):
        return "n/a"


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------


def slide_title(prs: Presentation) -> None:
    slide = blank(prs)
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY
    panel.line.fill.background()
    panel.shadow.inherit = False

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.05), Inches(1.5), Inches(0.07))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    accent.shadow.inherit = False

    label(
        slide,
        1.0,
        2.0,
        11.0,
        1.0,
        "Agentic Invoice-to-Payment Automation",
        size=44,
        color=WHITE,
        bold=True,
    )
    label(
        slide,
        1.0,
        3.35,
        11.0,
        0.8,
        "Email to posted payment journal, with a deterministic match and a human in the loop",
        size=19,
        color=RGBColor(0xC5, 0xD3, 0xE0),
    )
    label(slide, 1.0, 5.4, 11.0, 0.4, "Ashutosh Sharma", size=16, color=WHITE, bold=True)
    label(
        slide,
        1.0,
        5.78,
        11.0,
        0.4,
        "AI Engineer  |  Technical Assignment",
        size=13,
        color=RGBColor(0x9A, 0xAF, 0xC2),
    )


def slide_problem(prs: Presentation) -> None:
    slide = content(prs, "The problem", "What accounts payable looks like without automation")
    bullets(
        slide,
        0.7,
        1.85,
        6.1,
        [
            "Invoices arrive as email attachments: native PDFs, scans and HTML, in no "
            "fixed layout.",
            "A clerk keys header and line data into the ERP by hand.",
            "Each invoice is compared line by line against the purchase order and goods receipt.",
            "Variances are chased over email, and the reasoning is never written down.",
            "Duplicate invoices get paid twice because the check is a human memory.",
        ],
    )
    box(
        slide,
        7.3,
        1.85,
        5.3,
        1.15,
        "Per invoice, manual handling costs roughly 10-15x an automated touchless one",
        fill=PAPER,
        font_size=14,
        font_color=NAVY,
        bold=True,
    )
    box(
        slide,
        7.3,
        3.15,
        5.3,
        1.15,
        "Duplicate and overpayment leakage typically runs 0.5-1% of AP spend",
        fill=PAPER,
        font_size=14,
        font_color=NAVY,
        bold=True,
    )
    box(
        slide,
        7.3,
        4.45,
        5.3,
        1.15,
        "Auditors need the reason a payment was approved, not just the payment",
        fill=PAPER,
        font_size=14,
        font_color=NAVY,
        bold=True,
    )


def slide_solution(prs: Presentation) -> None:
    slide = content(prs, "What the agent does", "Eight steps, from unread email to posted journal")
    steps = [
        ("1", "Ingest", "Graph / Gmail /\nlocal folder"),
        ("2", "Parse", "Docling, OCR +\ntable structure"),
        ("3", "Extract", "Ollama into a\nPydantic schema"),
        ("4", "De-duplicate", "Vendor + number\n+ gross amount"),
        ("5", "Resolve vendor", "Fuzzy + pgvector\nvendor master"),
        ("6", "Fetch PO / GR", "ERP purchase order\nand goods receipt"),
        ("7", "Match", "2-way / 3-way,\nDecimal arithmetic"),
        ("8", "Post or escalate", "Journal, or a\nhuman approver"),
    ]
    left = 0.55
    for index, (number, name, detail) in enumerate(steps):
        x = left + (index % 4) * 3.12
        y = 1.95 + (index // 4) * 2.35
        box(slide, x, y, 2.85, 1.85, "", fill=PAPER, line=BORDER)
        label(slide, x + 0.16, y + 0.12, 0.5, 0.35, number, size=20, color=TEAL, bold=True)
        label(slide, x + 0.16, y + 0.55, 2.5, 0.4, name, size=15, color=NAVY, bold=True)
        label(slide, x + 0.16, y + 0.98, 2.5, 0.8, detail, size=11, color=MUTED)
        if index % 4 != 3:
            arrow(slide, x + 2.88, y + 0.92, x + 3.1, y + 0.92)


def slide_architecture(prs: Presentation) -> None:
    slide = content(prs, "Architecture", "One FastAPI service, four backing systems")

    box(slide, 0.55, 1.9, 2.35, 2.5, "", fill=PAPER, line=BORDER)
    label(slide, 0.62, 1.98, 2.2, 0.3, "SOURCES", size=10, color=MUTED, bold=True)
    for index, name in enumerate(["Microsoft Graph", "Gmail API", "Local folder", "REST upload"]):
        box(slide, 0.72, 2.35 + index * 0.48, 2.0, 0.4, name, fill=WHITE, font_size=11)

    box(slide, 3.25, 1.9, 6.0, 3.95, "", fill=RGBColor(0xEC, 0xF6, 0xF5), line=TEAL)
    label(
        slide,
        3.35,
        1.97,
        5.8,
        0.3,
        "INVOICE-TO-PAYMENT AGENT  (FastAPI)",
        size=10,
        color=TEAL,
        bold=True,
    )
    inner = [
        ("Ingestion", "mailbox providers"),
        ("Docling parser", "OCR + tables"),
        ("Extraction", "Ollama + schema"),
        ("LangGraph", "checkpointed workflow"),
        ("Matching engine", "Decimal, deterministic"),
        ("LlamaIndex RAG", "vendor + policy"),
    ]
    for index, (name, detail) in enumerate(inner):
        x = 3.42 + (index % 2) * 2.92
        y = 2.35 + (index // 2) * 1.14
        box(slide, x, y, 2.75, 0.98, "", fill=WHITE, line=BORDER)
        label(slide, x + 0.12, y + 0.12, 2.5, 0.35, name, size=12, color=NAVY, bold=True)
        label(slide, x + 0.12, y + 0.5, 2.5, 0.35, detail, size=10, color=MUTED)

    box(slide, 9.65, 1.9, 3.1, 3.95, "", fill=PAPER, line=BORDER)
    label(slide, 9.75, 1.97, 2.9, 0.3, "STATE AND EXTERNAL", size=10, color=MUTED, bold=True)
    ext = [
        ("PostgreSQL + pgvector", "runs, matches, journals,\naudit, embeddings, checkpoints"),
        ("ERP API", "mock SAP S/4HANA"),
        ("Ollama", "llama3.1:8b, nomic-embed"),
        ("Arize Phoenix", "OTLP traces"),
    ]
    for index, (name, detail) in enumerate(ext):
        y = 2.35 + index * 0.86
        box(slide, 9.82, y, 2.78, 0.74, "", fill=WHITE, line=BORDER)
        label(slide, 9.92, y + 0.05, 2.6, 0.3, name, size=11, color=NAVY, bold=True)
        label(slide, 9.92, y + 0.32, 2.6, 0.4, detail, size=8, color=MUTED)

    arrow(slide, 2.93, 3.15, 3.22, 3.15)
    arrow(slide, 9.28, 3.85, 9.62, 3.85)
    box(
        slide,
        3.25,
        6.05,
        6.0,
        0.62,
        "Human approver: reviews exceptions, approves or rejects, resumes the checkpointed run",
        fill=RGBColor(0xFD, 0xF3, 0xE4),
        line=AMBER,
        font_size=12,
        font_color=NAVY,
    )
    arrow(slide, 6.25, 5.88, 6.25, 6.02, AMBER)


def slide_graph(prs: Presentation) -> None:
    slide = content(
        prs,
        "The accounts-payable workflow",
        "A LangGraph state machine, checkpointed to Postgres at every node",
    )
    nodes = [
        "parse",
        "extract",
        "persist",
        "duplicate\ncheck",
        "resolve\nvendor",
        "fetch PO\n+ GR",
        "match",
    ]
    for index, name in enumerate(nodes):
        x = 0.55 + index * 1.79
        box(slide, x, 2.15, 1.6, 0.85, name, fill=PAPER, line=BORDER, font_size=11)
        if index < len(nodes) - 1:
            arrow(slide, x + 1.62, 2.57, x + 1.77, 2.57)

    box(
        slide,
        3.3,
        3.85,
        2.6,
        0.85,
        "post journal",
        fill=RGBColor(0xE7, 0xF5, 0xEC),
        line=RGBColor(0x2E, 0x8B, 0x57),
        font_size=13,
        font_color=RGBColor(0x1E, 0x5E, 0x3A),
        bold=True,
    )
    box(
        slide,
        7.3,
        3.85,
        2.9,
        0.85,
        "raise exception\n(graph interrupts)",
        fill=RGBColor(0xFD, 0xF3, 0xE4),
        line=AMBER,
        font_size=12,
        font_color=NAVY,
        bold=True,
    )
    box(slide, 7.3, 5.15, 2.9, 0.8, "human decision", fill=WHITE, line=AMBER, font_size=12)

    arrow(slide, 11.5, 3.0, 5.0, 3.8)
    arrow(slide, 11.9, 3.0, 8.6, 3.8, AMBER)
    arrow(slide, 8.75, 4.72, 8.75, 5.12, AMBER)
    arrow(slide, 7.28, 5.55, 4.6, 4.75, AMBER)

    label(slide, 4.05, 3.42, 3.0, 0.3, "in tolerance,\nunder ceiling", size=10, color=MUTED)
    label(
        slide, 9.15, 3.42, 3.2, 0.3, "variance / duplicate /\nlow confidence", size=10, color=MUTED
    )

    box(
        slide,
        0.55,
        6.25,
        12.2,
        0.75,
        "The interrupt persists full run state. The service can restart, redeploy, or wait three "
        "days for an approver, and the run resumes at exactly that node.",
        fill=PAPER,
        line=BORDER,
        font_size=12,
        font_color=NAVY,
    )


def slide_extraction(prs: Presentation) -> None:
    slide = content(
        prs,
        "Document processing and extraction",
        "Docling for layout, a local model for fields, schema validation for trust",
    )
    bullets(
        slide,
        0.7,
        1.85,
        6.2,
        [
            "Docling converts PDF, image, HTML and Office documents to Markdown with OCR and "
            "table-structure recognition enabled.",
            "Scanned invoices travel the same path as native PDFs, which matters because scans are "
            "the majority of real AP volume.",
            "A PyMuPDF text-layer fallback keeps the pipeline alive if Docling cannot "
            "handle a file.",
            "The model fills a strict Pydantic schema. A malformed response is repaired "
            "and retried "
            "rather than accepted.",
            "Every field carries its own confidence; the aggregate gates straight-through posting.",
        ],
        size=14,
        gap=0.72,
    )

    box(
        slide,
        7.3,
        1.85,
        5.4,
        0.5,
        "Extracted contract",
        fill=NAVY,
        line=NAVY,
        font_size=13,
        font_color=WHITE,
        bold=True,
    )
    fields = [
        "invoice_number, invoice_date, due_date",
        "purchase_order_number",
        "vendor_name, vendor_tax_id, vendor_iban",
        "currency, subtotal, tax_amount, total_amount",
        "payment_terms",
        "line_items[]: description, material_code,\nquantity, unit_price, line_total",
        "confidence, field_confidence{}",
    ]
    y = 2.45
    for field in fields:
        height = 0.62 if "\n" in field else 0.42
        box(slide, 7.3, y, 5.4, height, field, fill=PAPER, line=BORDER, font_size=11)
        y += height + 0.09


def slide_matching(prs: Presentation) -> None:
    slide = content(
        prs,
        "The matching engine",
        "The model extracts; deterministic code decides whether money moves",
    )
    rows = [
        ("Unit price", "2%", "percentage on the unit price"),
        ("Unit price, absolute", "5.00", "on the extended line value"),
        ("Quantity", "exact", "no tolerance for goods"),
        ("Invoice gross total", "1% or 10.00", "whichever the variance satisfies"),
        ("Tax", "1.00", "absorbs rounding"),
        ("Auto-post ceiling", "25,000.00", "segregation of duties"),
        ("Min extraction confidence", "0.75", "below this, a human looks"),
        ("Vendor name similarity", "88", "fuzzy against master + aliases"),
    ]
    box(slide, 0.55, 1.85, 7.1, 0.42, "", fill=NAVY, line=NAVY)
    label(slide, 0.7, 1.89, 3.0, 0.34, "RULE", size=11, color=WHITE, bold=True)
    label(slide, 3.9, 1.89, 1.6, 0.34, "DEFAULT", size=11, color=WHITE, bold=True)
    label(slide, 5.6, 1.89, 2.0, 0.34, "APPLIED TO", size=11, color=WHITE, bold=True)
    for index, (rule, default, note) in enumerate(rows):
        y = 2.32 + index * 0.44
        shade = PAPER if index % 2 == 0 else WHITE
        box(slide, 0.55, y, 7.1, 0.44, "", fill=shade, line=BORDER, shape=MSO_SHAPE.RECTANGLE)
        label(slide, 0.7, y + 0.06, 3.1, 0.34, rule, size=11, color=SLATE)
        label(slide, 3.9, y + 0.06, 1.6, 0.34, default, size=11, color=TEAL, bold=True)
        label(slide, 5.6, y + 0.06, 2.0, 0.34, note, size=10, color=MUTED)

    box(
        slide,
        8.0,
        1.85,
        4.75,
        2.0,
        "The absolute price tolerance governs the extended line value, not the unit price.\n\n"
        "A flat 5.00 per-unit floor waves through a 34% overcharge on a 14.50 part, and "
        "hundreds of euros once multiplied across quantity.",
        fill=RGBColor(0xFD, 0xF3, 0xE4),
        line=AMBER,
        font_size=12,
        font_color=NAVY,
    )
    box(
        slide,
        8.0,
        4.05,
        4.75,
        1.75,
        "Line pairing falls back in three steps: PO line number, then material code, then fuzzy "
        "description. Real invoices carry neither of the first two, and hard-joining "
        "manufactures exceptions a human then clears by hand.",
        fill=PAPER,
        line=BORDER,
        font_size=12,
        font_color=NAVY,
    )


def slide_exceptions(prs: Presentation) -> None:
    slide = content(
        prs,
        "Exceptions and the human in the loop",
        "Fourteen typed exceptions, each with the policy text that produced it",
    )
    types = [
        "price_variance",
        "quantity_variance",
        "total_variance",
        "tax_variance",
        "missing_po",
        "missing_goods_receipt",
        "duplicate_invoice",
        "unknown_vendor",
        "low_confidence_extraction",
        "line_not_on_po",
        "over_auto_post_ceiling",
        "currency_mismatch",
        "unapplied_remittance",
        "erp_posting_failed",
    ]
    for index, name in enumerate(types):
        x = 0.55 + (index % 4) * 3.12
        y = 1.9 + (index // 4) * 0.55
        box(slide, x, y, 2.95, 0.45, name, fill=PAPER, line=BORDER, font_size=10)

    box(
        slide,
        0.55,
        4.15,
        5.95,
        1.5,
        "Each case carries a severity, a plain-language summary, a suggested action, and the "
        "variance detail. The relevant AP policy is retrieved from the knowledge base and "
        "attached, so the approver sees the rule, not only the number that broke it.",
        fill=WHITE,
        line=BORDER,
        font_size=12,
        font_color=SLATE,
    )
    box(
        slide,
        6.8,
        4.15,
        5.95,
        1.5,
        "Duplicate detection keys on vendor, invoice number and gross amount. Invoice date is "
        "deliberately excluded: vendors re-issue the same invoice with a fresh print date, and "
        "including it would let a genuine duplicate through.",
        fill=RGBColor(0xFD, 0xF3, 0xE4),
        line=AMBER,
        font_size=12,
        font_color=NAVY,
    )
    box(
        slide,
        0.55,
        5.9,
        12.2,
        0.75,
        "GET /exceptions  ->  human reviews  ->  POST /exceptions/{id}/decision  ->  the "
        "checkpointed run resumes and posts, or closes as rejected. Either way the decision, "
        "the approver and the note land in the audit trail.",
        fill=NAVY,
        line=NAVY,
        font_size=12,
        font_color=WHITE,
    )


def slide_ar(prs: Presentation) -> None:
    slide = content(
        prs, "The accounts-receivable mirror", "The same pattern, pointed at incoming remittances"
    )
    nodes = ["parse", "extract\nremittance", "fetch open\nAR items", "apply cash"]
    for index, name in enumerate(nodes):
        x = 1.4 + index * 2.85
        box(slide, x, 2.2, 2.35, 1.0, name, fill=PAPER, line=BORDER, font_size=13)
        if index < len(nodes) - 1:
            arrow(slide, x + 2.38, 2.7, x + 2.82, 2.7)
    bullets(
        slide,
        0.7,
        4.0,
        12.0,
        [
            "A customer remittance advice lists which invoices a payment covers, and for how much.",
            "Open AR items are pulled from the ERP and matched against the advice line by line.",
            "Cash is applied per item; a residual outside the remittance tolerance stays open.",
            "An unmatched payment raises unapplied_remittance rather than being force-applied.",
            "It reuses the same parser, the same extraction path and the same audit trail.",
        ],
        size=14,
        gap=0.5,
    )


def slide_audit(prs: Presentation) -> None:
    slide = content(prs, "Audit trail", "Why a payment was made, reconstructable without a model")
    events = [
        "email_received",
        "document_parsed",
        "fields_extracted",
        "duplicate_checked",
        "vendor_resolved",
        "po_fetched",
        "goods_receipt_fetched",
        "match_evaluated",
        "exception_raised",
        "human_decision",
        "journal_posted",
        "cash_applied",
        "run_failed",
    ]
    for index, name in enumerate(events):
        y = 1.9 + index * 0.36
        box(
            slide,
            0.55,
            y,
            3.4,
            0.3,
            name,
            fill=PAPER,
            line=BORDER,
            font_size=10,
            shape=MSO_SHAPE.RECTANGLE,
        )

    box(
        slide,
        4.35,
        1.9,
        8.4,
        2.15,
        "Every event records the acting node, a summary, a structured payload, elapsed "
        "milliseconds, and - critically - the source email message id and the SHA-256 of the "
        "source document.\n\nThat hash is what ties a posted payment back to the exact bytes "
        "that justified it.",
        fill=WHITE,
        line=TEAL,
        font_size=13,
        font_color=SLATE,
    )
    box(
        slide,
        4.35,
        4.25,
        8.4,
        1.15,
        "Filtered by run_id, events return in workflow sequence - the order an auditor reads "
        "them. Unfiltered, newest first - the order an operator reads them.",
        fill=PAPER,
        line=BORDER,
        font_size=12,
        font_color=SLATE,
    )
    box(
        slide,
        4.35,
        5.6,
        8.4,
        1.0,
        "The straight-through rate is derived from run rows, not counted into a separate tally, "
        "so the headline metric cannot drift away from what the audit events say happened.",
        fill=NAVY,
        line=NAVY,
        font_size=12,
        font_color=WHITE,
    )


def slide_observability(prs: Presentation) -> None:
    slide = content(prs, "Observability", "Arize Phoenix over OpenInference, and never fatal")
    bullets(
        slide,
        0.7,
        1.9,
        6.2,
        [
            "LangChain and LlamaIndex instrumentors export spans over OTLP to Phoenix.",
            "Every extraction, retrieval and graph node appears as a span with latency and tokens.",
            "A slow or failing invoice is inspected as a trace, not reconstructed from log lines.",
            "Tracing setup is wrapped: an unreachable collector logs a warning and the service "
            "runs on. Observability is not the product.",
            "Health reports each dependency separately, so degraded is distinguishable from down.",
        ],
        size=14,
        gap=0.66,
    )

    box(
        slide,
        7.3,
        1.9,
        5.45,
        0.5,
        "GET /api/v1/health",
        fill=NAVY,
        line=NAVY,
        font_size=13,
        font_color=WHITE,
        bold=True,
    )
    checks = [
        ("postgres", "service cannot run without it", RGBColor(0xC0, 0x39, 0x2B)),
        ("ollama", "no new processing, queue readable", AMBER),
        ("erp", "no new posting, queue readable", AMBER),
        ("phoenix", "traces lost only", MUTED),
    ]
    for index, (name, note, colour) in enumerate(checks):
        y = 2.55 + index * 0.72
        box(slide, 7.3, y, 5.45, 0.62, "", fill=PAPER, line=BORDER)
        label(slide, 7.48, y + 0.06, 2.0, 0.3, name, size=12, color=NAVY, bold=True)
        label(slide, 7.48, y + 0.32, 5.1, 0.28, note, size=10, color=colour)

    box(
        slide,
        7.3,
        5.5,
        5.45,
        0.95,
        "Postgres down -> down.\nAnything else degraded -> degraded, because the audit trail and "
        "approval queue stay readable.",
        fill=RGBColor(0xEC, 0xF6, 0xF5),
        line=TEAL,
        font_size=12,
        font_color=NAVY,
    )


def slide_evaluation(prs: Presentation, metrics: dict[str, Any]) -> None:
    slide = content(prs, "Evaluation", "Scored against a hand-written answer key, not a vibe check")

    extraction = metrics.get("extraction", {})
    matching = metrics.get("matching", {})
    stp = metrics.get("straight_through", {})
    processed = metrics.get("documents_processed", 0)
    total = metrics.get("documents_total", 0)

    # Matching and STP are measured only in e2e mode. Rendering an unmeasured metric as
    # "0%" reads as a failure on a slide, so an empty denominator shows as "--".
    measured = bool(matching.get("outcomes_total"))
    tiles = [
        (pct(extraction.get("field_accuracy")), "Extraction\nfield accuracy", TEAL),
        (pct(extraction.get("line_item_accuracy")), "Line item\naccuracy", TEAL),
        (pct(matching.get("match_rate")) if measured else "--", "Match\nrate", NAVY),
        (
            pct(matching.get("exception_detection_rate")) if measured else "--",
            "Exception\ndetection rate",
            AMBER,
        ),
        (pct(stp.get("stp_rate")) if measured else "--", "Straight-through\nrate", NAVY),
    ]
    for index, (value, caption, colour) in enumerate(tiles):
        metric_tile(slide, 0.55 + index * 2.5, 1.85, 2.3, 1.55, value, caption, colour)

    caption = (
        f"{processed} of {total} sample documents scored, "
        f"{metrics.get('documents_failed', 0)} failed  |  "
        f"mode: {metrics.get('mode', 'extraction')}  |  model: llama3.1:8b via Ollama"
        if metrics
        else "Run: python -m evaluation.run_evaluation --mode extraction"
    )
    if metrics and not measured:
        caption += "    (--  measured only in e2e mode, against the running stack)"
    label(slide, 0.6, 3.55, 12.2, 0.35, caption, size=11, color=MUTED)
    bullets(
        slide,
        0.7,
        4.05,
        6.2,
        [
            "Seven documents covering a clean three-way match, a clean two-way match, a price "
            "variance, a goods-receipt shortfall, an unknown vendor with no PO, a duplicate "
            "resubmission, and an invoice above the auto-post ceiling.",
            "Per-field accuracy is reported field by field, so a systematically weak field is "
            "visible rather than averaged away.",
        ],
        size=13,
        gap=1.15,
    )

    box(
        slide,
        7.3,
        4.05,
        5.45,
        2.3,
        "Exception detection rate is the number that matters most.\n\n"
        "A missed variance is an incorrect payment. A false exception only costs a human a "
        "minute. The tolerances are tuned accordingly - the engine is built to over-escalate "
        "rather than over-pay.",
        fill=RGBColor(0xFD, 0xF3, 0xE4),
        line=AMBER,
        font_size=13,
        font_color=NAVY,
    )


def slide_deployment(prs: Presentation) -> None:
    slide = content(
        prs, "Deployment", "One command, seven services, ordered by health not by sleep"
    )
    services = [
        ("postgres", "pgvector/pgvector:pg16", "state + embeddings + checkpoints"),
        ("ollama", "ollama/ollama", "llama3.1:8b, nomic-embed-text"),
        ("ollama-init", "one-shot", "pulls models, then exits"),
        ("phoenix", "arizephoenix/phoenix", "trace collector and UI"),
        ("mock-erp", "built from repo", "POs, GRs, vendors, journals, AR"),
        ("migrate", "one-shot", "alembic upgrade head"),
        ("api", "built from repo", "the agent, port 8000"),
    ]
    box(slide, 0.55, 1.85, 7.4, 0.42, "", fill=NAVY, line=NAVY)
    label(slide, 0.7, 1.89, 2.0, 0.34, "SERVICE", size=11, color=WHITE, bold=True)
    label(slide, 2.8, 1.89, 2.4, 0.34, "IMAGE", size=11, color=WHITE, bold=True)
    label(slide, 5.3, 1.89, 2.5, 0.34, "ROLE", size=11, color=WHITE, bold=True)
    for index, (name, image, role) in enumerate(services):
        y = 2.32 + index * 0.46
        shade = PAPER if index % 2 == 0 else WHITE
        box(slide, 0.55, y, 7.4, 0.46, "", fill=shade, line=BORDER, shape=MSO_SHAPE.RECTANGLE)
        label(slide, 0.7, y + 0.07, 2.1, 0.34, name, size=11, color=NAVY, bold=True)
        label(slide, 2.8, y + 0.07, 2.5, 0.34, image, size=10, color=MUTED)
        label(slide, 5.3, y + 0.07, 2.6, 0.34, role, size=10, color=SLATE)

    box(
        slide,
        8.3,
        1.85,
        4.45,
        1.25,
        "cp .env.example .env\ndocker compose up -d --build",
        fill=NAVY,
        line=NAVY,
        font_size=13,
        font_color=WHITE,
        bold=True,
    )
    box(
        slide,
        8.3,
        3.25,
        4.45,
        2.3,
        "Startup order is enforced with health conditions, not sleeps.\n\n"
        "Postgres must accept connections before migrations run. Migrations must finish before "
        "the API starts. The models must be pulled before the first extraction is attempted.",
        fill=PAPER,
        line=BORDER,
        font_size=12,
        font_color=SLATE,
    )
    box(
        slide,
        0.55,
        5.75,
        12.2,
        0.72,
        "Swagger UI at :8000/docs   |   Mock ERP at :8081/docs   |   Phoenix at :6006   |   "
        "OpenAPI committed at docs/openapi.json",
        fill=RGBColor(0xEC, 0xF6, 0xF5),
        line=TEAL,
        font_size=12,
        font_color=NAVY,
    )


def slide_decisions(prs: Presentation) -> None:
    slide = content(prs, "Design decisions", "The four that shaped everything else")
    cards = [
        (
            "The LLM extracts; code decides",
            "Nothing about whether to pay money is left to a language model. The model turns a "
            "document into fields; Decimal arithmetic compares them to the PO. Every posting "
            "decision is reproducible, and tightening a tolerance is configuration, not prompting.",
        ),
        (
            "Human-in-the-loop by checkpointed interrupt",
            "The approval pause is a LangGraph interrupt over a Postgres checkpointer. Full run "
            "state persists, so the service can be redeployed while an invoice waits, and the "
            "resumed run continues from exactly that node.",
        ),
        (
            "Two vector indexes, not one",
            "Vendor master and AP policy live in separate pgvector tables because they answer "
            "different questions. Blended into a single top-k, policy prose consistently outranks "
            "the vendor row the query actually needed.",
        ),
        (
            "Metrics derived, not asserted",
            "The straight-through rate is computed from run rows rather than counted into a "
            "separate tally, so the headline number cannot drift away from what the audit trail "
            "says actually happened.",
        ),
    ]
    for index, (title, body) in enumerate(cards):
        x = 0.55 + (index % 2) * 6.35
        y = 1.9 + (index // 2) * 2.35
        box(slide, x, y, 6.0, 2.1, "", fill=PAPER, line=BORDER)
        label(slide, x + 0.22, y + 0.16, 5.6, 0.4, title, size=15, color=NAVY, bold=True)
        label(slide, x + 0.22, y + 0.64, 5.6, 1.35, body, size=11, color=SLATE)


def slide_roadmap(prs: Presentation) -> None:
    slide = content(
        prs,
        "Limitations and what production needs",
        "Stated plainly, because a prototype that hides its edges is worse than useless",
    )
    left = [
        (
            "The ERP is a mock",
            "mock_erp mirrors the S/4HANA OData shape and reseeds on startup "
            "for reproducible evaluation. A real integration swaps "
            "erp/client.py; nothing upstream changes.",
        ),
        (
            "Extraction tracks the local model",
            "llama3.1:8b keeps the whole stack on one machine "
            "with no external API calls - the right default "
            "for finance documents. A larger model measurably "
            "improves dense line-item extraction.",
        ),
    ]
    right = [
        (
            "Polling is request-triggered",
            "POST /mailbox/poll drives a poll. Production would run "
            "it on a scheduler, or move to Graph change "
            "notifications instead of interval polling.",
        ),
        (
            "No currency conversion",
            "A currency mismatch between invoice and PO is detected and "
            "escalated rather than converted at a daily rate.",
        ),
    ]
    for column, items in enumerate([left, right]):
        for index, (title, body) in enumerate(items):
            x = 0.55 + column * 6.35
            y = 1.9 + index * 1.95
            box(slide, x, y, 6.0, 1.7, "", fill=WHITE, line=BORDER)
            label(slide, x + 0.22, y + 0.14, 5.6, 0.35, title, size=14, color=NAVY, bold=True)
            label(slide, x + 0.22, y + 0.55, 5.6, 1.05, body, size=11, color=SLATE)

    box(
        slide,
        0.55,
        5.72,
        12.2,
        0.82,
        "Next: Graph change notifications instead of polling  |  vendor-specific extraction "
        "templates for high-volume suppliers  |  a learned tolerance recommendation from "
        "historical approver decisions  |  RAGAs scoring of the policy-retrieval path",
        fill=NAVY,
        line=NAVY,
        font_size=12,
        font_color=WHITE,
    )
    label(
        slide,
        0.6,
        6.68,
        12.1,
        0.35,
        "Ashutosh Sharma  |  Agentic Invoice-to-Payment Automation  |  Thank you",
        size=11,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def build(output: Path) -> Path:
    metrics = load_metrics()
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_architecture(prs)
    slide_graph(prs)
    slide_extraction(prs)
    slide_matching(prs)
    slide_exceptions(prs)
    slide_ar(prs)
    slide_audit(prs)
    slide_observability(prs)
    slide_evaluation(prs, metrics)
    slide_deployment(prs)
    slide_decisions(prs)
    slide_roadmap(prs)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return output


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "-o",
        "--output",
        type=Path,
        default=ROOT / "docs" / "Invoice-to-Payment-Agent-Solution-Deck.pptx",
    )
    args = parser.parse_args()
    path = build(args.output)
    print(f"Deck written to {path.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
